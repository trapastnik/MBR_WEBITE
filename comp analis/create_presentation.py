#!/usr/bin/env python3
"""Generate company profile presentations for Danyelza/Y-mAbs and Unituxin/United Therapeutics."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# === COLORS ===
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)

# Danyelza / Y-mAbs palette
DZ_PRIMARY = RGBColor(0x1B, 0x3A, 0x5C)      # dark navy
DZ_ACCENT = RGBColor(0x3A, 0x8F, 0xD6)       # bright blue
DZ_ACCENT_LINE = RGBColor(0x5B, 0xB5, 0xF0)  # light blue accent
DZ_TABLE_HDR = RGBColor(0x1B, 0x3A, 0x5C)
DZ_TABLE_ALT = RGBColor(0xE3, 0xEE, 0xF8)
DZ_LIGHT_BG = RGBColor(0xF5, 0xF8, 0xFC)
DZ_SUBTITLE = RGBColor(0xA0, 0xC4, 0xE8)

# Unituxin / United Therapeutics palette
UT_PRIMARY = RGBColor(0x2B, 0x4C, 0x3F)      # dark teal-green
UT_ACCENT = RGBColor(0x4A, 0x8C, 0x72)       # muted green
UT_ACCENT_LINE = RGBColor(0x6B, 0xB3, 0x95)  # soft sage accent
UT_TABLE_HDR = RGBColor(0x2B, 0x4C, 0x3F)
UT_TABLE_ALT = RGBColor(0xE8, 0xF2, 0xED)
UT_LIGHT_BG = RGBColor(0xF2, 0xF8, 0xF5)
UT_SUBTITLE = RGBColor(0xA3, 0xCC, 0xB8)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = Inches(13.333)


def add_background(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK_GRAY,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullets(slide, left, top, width, height, items, font_size=14, color=DARK_GRAY, spacing=8):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(spacing)
    return txBox


def add_table(slide, left, top, width, rows_data, col_widths=None, header_bg=DZ_TABLE_HDR, alt_bg=DZ_TABLE_ALT):
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(0.4 * n_rows))
    table = tbl_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r, row in enumerate(rows_data):
        for c, txt in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(txt)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(12) if r > 0 else Pt(13)
                para.font.name = "Calibri"
                if r == 0:
                    para.font.bold = True
                    para.font.color.rgb = WHITE
                    para.alignment = PP_ALIGN.CENTER
                else:
                    para.font.color.rgb = DARK_GRAY
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = alt_bg
    return tbl_shape


def section_title_bar(slide, title, subtitle, primary_color, accent_color, subtitle_color):
    """Standard title bar for content slides — shows company branding."""
    add_background(slide, WHITE)
    rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), primary_color)
    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.7),
                 title, font_size=28, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.6), Inches(0.75), Inches(12), Inches(0.4),
                     subtitle, font_size=14, color=subtitle_color)
    rect(slide, Inches(0), Inches(1.2), SLIDE_W, Inches(0.06), accent_color)


# ================================================================
# SLIDE 1: GENERAL TITLE
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, RGBColor(0x0F, 0x1F, 0x33))
rect(slide, Inches(2), Inches(2.0), Inches(9.333), Inches(0.05), DZ_ACCENT)
add_text_box(slide, Inches(1), Inches(2.3), Inches(11.333), Inches(1.5),
             "Anti-GD2 Monoclonal Antibodies\nin Neuroblastoma",
             font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.2), Inches(11.333), Inches(0.6),
             "Company Profiles: Danyelza (Y-mAbs / SERB)  &  Unituxin (United Therapeutics)",
             font_size=20, color=RGBColor(0xAA, 0xCC, 0xEE), alignment=PP_ALIGN.CENTER)
rect(slide, Inches(2), Inches(5.1), Inches(9.333), Inches(0.05), DZ_ACCENT)


# ================================================================
#  PART 1 — DANYELZA / Y-mAbs  (slides 2–7)
# ================================================================

# --- SLIDE 2: Danyelza section cover ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DZ_PRIMARY)
rect(slide, Inches(0.8), Inches(2.8), Inches(11.733), Inches(0.05), DZ_ACCENT_LINE)
add_text_box(slide, Inches(1), Inches(3.1), Inches(11.333), Inches(1.0),
             "DANYELZA (Naxitamab-gqgk)",
             font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.3), Inches(11.333), Inches(0.5),
             "Y-mAbs Therapeutics  /  SERB Pharmaceuticals",
             font_size=22, color=DZ_SUBTITLE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.0), Inches(11.333), Inches(0.5),
             "Humanized anti-GD2 monoclonal antibody (hu3F8) for relapsed/refractory high-risk neuroblastoma",
             font_size=15, color=RGBColor(0x88, 0xAA, 0xCC), alignment=PP_ALIGN.CENTER)
rect(slide, Inches(0.8), Inches(5.6), Inches(11.733), Inches(0.05), DZ_ACCENT_LINE)


# --- SLIDE 3: Danyelza — Overview ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_title_bar(slide, "DANYELZA  |  Overview", None, DZ_PRIMARY, DZ_ACCENT, DZ_SUBTITLE)

data = [
    ["Parameter", "Details"],
    ["Generic Name", "Naxitamab-gqgk"],
    ["Brand Name", "Danyelza"],
    ["Type", "Humanized anti-GD2 monoclonal antibody (hu3F8)"],
    ["Target", "GD2 (disialoganglioside)"],
    ["Indication", "Relapsed/refractory high-risk neuroblastoma (in combination with GM-CSF)"],
    ["FDA Approval", "November 25, 2020 (accelerated approval)"],
    ["Original Developer", "Memorial Sloan Kettering Cancer Center (Dr. Nai-Kong Cheung)"],
    ["Commercializer", "Y-mAbs Therapeutics (2015–2025) → SERB Pharmaceuticals (acquired Sept 2025)"],
    ["Developed In-House?", "NO — Academic (MSK) developed, licensed to Y-mAbs"],
    ["FDA Designations", "Accelerated Approval, Breakthrough Therapy, Priority Review,\nOrphan Drug, Rare Pediatric Disease, PRV (sold)"],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), data,
          col_widths=[Inches(2.8), Inches(9.5)], header_bg=DZ_TABLE_HDR, alt_bg=DZ_TABLE_ALT)


# --- SLIDE 4: Danyelza — Development History ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_title_bar(slide, "DANYELZA  |  Development History",
                  "30+ years of translational research at Memorial Sloan Kettering", DZ_PRIMARY, DZ_ACCENT, DZ_SUBTITLE)

data = [
    ["Period", "Milestone"],
    ["1980s", "Dr. Nai-Kong Cheung develops murine anti-GD2 antibody 3F8 at MSK"],
    ["Mid-1980s", "Clinical trials of 3F8 begin at MSK"],
    ["1990s–2000s", "Extensive clinical experience with 3F8 in hundreds of neuroblastoma patients"],
    ["2000s–2010s", "Development of humanized version hu3F8 (naxitamab) to reduce immunogenicity"],
    ["2015", "Y-mAbs Therapeutics founded to commercialize MSK anti-GD2 portfolio"],
    ["Sept 2018", "Y-mAbs IPO on Nasdaq (~$86M raised at ~$14/share)"],
    ["Nov 25, 2020", "FDA accelerated approval of Danyelza"],
    ["Aug 2025", "SERB Pharmaceuticals agrees to acquire Y-mAbs"],
    ["Sept 16, 2025", "SERB acquisition closes; Y-mAbs delisted from Nasdaq"],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), data,
          col_widths=[Inches(2.0), Inches(10.3)], header_bg=DZ_TABLE_HDR, alt_bg=DZ_TABLE_ALT)

add_bullets(slide, Inches(0.5), Inches(6.0), Inches(12), Inches(1.2), [
    "Pivotal Study 201 (hu3F8 + GM-CSF) conducted at MSK — demonstrated efficacy in relapsed/refractory patients",
    "3F8 (murine, never commercially approved) → hu3F8 (humanized, lower immunogenicity, more durable treatment)",
], font_size=13, color=MED_GRAY)


# --- SLIDE 5: Danyelza — Y-mAbs Company Profile ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_title_bar(slide, "DANYELZA  |  Y-mAbs Therapeutics — Company Profile",
                  "The commercializing company before SERB acquisition", DZ_PRIMARY, DZ_ACCENT, DZ_SUBTITLE)

data = [
    ["Parameter", "Details"],
    ["Founded", "2015 by Thomas Gad (Danish entrepreneur)"],
    ["Incorporated", "Delaware, USA"],
    ["Purpose", "Created specifically to commercialize Dr. Cheung's MSK anti-GD2 portfolio"],
    ["IPO", "September 28, 2018 — Nasdaq (YMAB), ~$14/share, ~$86M raised"],
    ["Stage at IPO", "Clinical-stage, pre-revenue, no approved products"],
    ["Manufacturing", "Via contract manufacturing organizations (CMOs)"],
    ["Key Assets", "Naxitamab (Danyelza), Omburtamab (8H9, anti-B7-H3), murine 3F8"],
    ["Value Proposition", "MSK-licensed assets + Dr. Cheung's 30+ year clinical track record"],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), data,
          col_widths=[Inches(2.8), Inches(9.5)], header_bg=DZ_TABLE_HDR, alt_bg=DZ_TABLE_ALT)

add_bullets(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(1.5), [
    "Classic academic spinout / biotech startup — pre-revenue, cash-burning, small team",
    "Entire company built around Danyelza as the core commercial asset",
    "License from MSK: exclusive worldwide rights to naxitamab, 3F8, and omburtamab",
    "IP: patents covering humanized 3F8 (composition of matter, methods of use); Dr. Cheung as inventor",
], font_size=13, color=MED_GRAY)


# --- SLIDE 6: Danyelza — Revenue & Commercial ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_title_bar(slide, "DANYELZA  |  Revenue & Commercial Performance",
                  "Strong growth trajectory since 2021 launch", DZ_PRIMARY, DZ_ACCENT, DZ_SUBTITLE)

data = [
    ["Year", "Net Revenue", "YoY Growth", "Notes"],
    ["2021", "$34.9M", "—", "Launch year (product: $32.9M + licensing: $2.0M)"],
    ["2022", "~$49.3M", "+41%", "Record Q4 at $16.4M (+31% sequential)"],
    ["2023", "$84.3M", "+71%", "Record year; strong international expansion"],
    ["2024", "$87.7M", "+4%", "International revenue $19.2M (+16% YoY)"],
    ["Q1 2025", "$20.9M", "+8% YoY", "—"],
    ["Q2 2025", "$19.5M", "Above guidance", "Exceeded $17–19M guidance range"],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), data,
          col_widths=[Inches(1.5), Inches(2.0), Inches(2.0), Inches(6.8)],
          header_bg=DZ_TABLE_HDR, alt_bg=DZ_TABLE_ALT)

add_bullets(slide, Inches(0.5), Inches(5.2), Inches(12), Inches(2.0), [
    "Revenue grew from ~$35M to ~$88M in 3 years (2021–2024)",
    "International expansion is a key growth driver (~$19.2M in 2024, +16% YoY)",
    "H1 2025 revenue: $40.4M — on track for continued growth under SERB ownership",
], font_size=14, color=DARK_GRAY)


# --- SLIDE 7: Danyelza — Valuation & Acquisition ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_title_bar(slide, "DANYELZA  |  Valuation & SERB Acquisition",
                  "From IPO to acquisition — the Y-mAbs trajectory", DZ_PRIMARY, DZ_ACCENT, DZ_SUBTITLE)

data = [
    ["Period", "Price Range", "Market Cap", "Event"],
    ["IPO (Sept 2018)", "~$14/share", "~$400–500M", "Clinical-stage launch on Nasdaq"],
    ["2019–2020", "Rising", "Growing", "Positive clinical data, BLA filing"],
    ["Peak (~2020)", "~$40–55+/share", "~$1.5–2.0B", "FDA approval euphoria"],
    ["2021–2023", "$3–10/share", "$100–400M", "Commercial disappointment, omburtamab CRL,\nbiotech downturn"],
    ["52-wk range (2024–25)", "$3.55–$16.11", "—", "Volatile trading"],
    ["Acquisition (Aug 2025)", "$8.60/share", "~$412M", "SERB buyout at 105% premium"],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), data,
          col_widths=[Inches(2.5), Inches(2.2), Inches(2.2), Inches(5.4)],
          header_bg=DZ_TABLE_HDR, alt_bg=DZ_TABLE_ALT)

add_bullets(slide, Inches(0.5), Inches(5.3), Inches(12), Inches(2.0), [
    "SERB acquisition: $8.60/share cash (~$412M total), 105% premium over pre-announcement price",
    "~16% of stockholders entered tender and support agreement",
    "Y-mAbs now operates as subsidiary of SERB S.A.S.; delisted from Nasdaq Sept 16, 2025",
    "Classic biotech boom-and-bust: peak ~$2B market cap → acquired at ~$412M",
], font_size=14, color=DARK_GRAY)


# ================================================================
#  PART 2 — UNITUXIN / United Therapeutics  (slides 8–13)
# ================================================================

# --- SLIDE 8: Unituxin section cover ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, UT_PRIMARY)
rect(slide, Inches(0.8), Inches(2.8), Inches(11.733), Inches(0.05), UT_ACCENT_LINE)
add_text_box(slide, Inches(1), Inches(3.1), Inches(11.333), Inches(1.0),
             "UNITUXIN (Dinutuximab)",
             font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.3), Inches(11.333), Inches(0.5),
             "United Therapeutics Corporation",
             font_size=22, color=UT_SUBTITLE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.0), Inches(11.333), Inches(0.5),
             "Chimeric anti-GD2 monoclonal antibody (ch14.18) for high-risk neuroblastoma",
             font_size=15, color=RGBColor(0x99, 0xBB, 0xAA), alignment=PP_ALIGN.CENTER)
rect(slide, Inches(0.8), Inches(5.6), Inches(11.733), Inches(0.05), UT_ACCENT_LINE)


# --- SLIDE 9: Unituxin — Overview ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_title_bar(slide, "UNITUXIN  |  Overview", None, UT_PRIMARY, UT_ACCENT, UT_SUBTITLE)

data = [
    ["Parameter", "Details"],
    ["Generic Name", "Dinutuximab"],
    ["Brand Name", "Unituxin"],
    ["Type", "Chimeric anti-GD2 monoclonal antibody (ch14.18)"],
    ["Target", "GD2 (disialoganglioside)"],
    ["Indication", "High-risk neuroblastoma (pediatric, frontline post-consolidation)"],
    ["Combination", "GM-CSF + IL-2 + isotretinoin"],
    ["FDA Approval", "March 10, 2015"],
    ["Current Owner", "United Therapeutics (U.S.); Recordati/EUSA Pharma (EU — Qarziba)"],
    ["Developed In-House?", "NO — Government/academic (NCI, COG, Scripps, UCSD)"],
    ["FDA Designations", "Priority Review, Breakthrough Therapy, Orphan Drug,\nRare Pediatric Disease, PRV ($67–350M historical value)"],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), data,
          col_widths=[Inches(2.8), Inches(9.5)], header_bg=UT_TABLE_HDR, alt_bg=UT_TABLE_ALT)


# --- SLIDE 10: Unituxin — Development History ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_title_bar(slide, "UNITUXIN  |  Development History",
                  "Government & academic-driven development through NCI and COG", UT_PRIMARY, UT_ACCENT, UT_SUBTITLE)

data = [
    ["Period", "Milestone"],
    ["1980s", "Dr. Ralph Reisfeld develops murine anti-GD2 antibody 14.18 at Scripps Research Institute"],
    ["Late 1980s–90s", "Chimeric ch14.18 engineered; Dr. Alice Yu (UCSD) advances clinical development"],
    ["2000s", "NCI funds research, manufactures antibody via Biologics Resources Branch, sponsors IND via CTEP"],
    ["2006–2010", "Pivotal ANBL0032 Phase III trial by COG across 200+ children's hospitals"],
    ["2010", "Yu et al. publish landmark results in NEJM — significant OS/EFS improvement"],
    ["March 10, 2015", "FDA approves dinutuximab as Unituxin"],
    ["2017", "United Therapeutics withdraws Unituxin from EU market (manufacturing difficulties)"],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), data,
          col_widths=[Inches(2.0), Inches(10.3)], header_bg=UT_TABLE_HDR, alt_bg=UT_TABLE_ALT)

add_bullets(slide, Inches(0.5), Inches(5.7), Inches(12), Inches(1.5), [
    "Pivotal ANBL0032: ch14.18 + GM-CSF + IL-2 + isotretinoin vs. isotretinoin alone — significant OS/EFS benefit",
    "Rare example of a drug developed almost entirely through government and academic effort, then licensed to a private company",
    "EU gap filled by Qarziba (dinutuximab beta) — Apeiron → EUSA Pharma → Recordati ($845M acquisition, 2021)",
], font_size=13, color=MED_GRAY)


# --- SLIDE 11: Unituxin — United Therapeutics Company Profile ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_title_bar(slide, "UNITUXIN  |  United Therapeutics — Company Profile",
                  "The commercializing company — an established PAH leader", UT_PRIMARY, UT_ACCENT, UT_SUBTITLE)

data = [
    ["Parameter", "Details"],
    ["Founded", "1996 by Martine Rothblatt"],
    ["Motivation", "Daughter diagnosed with pulmonary arterial hypertension (PAH)"],
    ["Core Business", "PAH treatments (Remodulin, Tyvaso, Orenitram, Adcirca)"],
    ["Annual Revenue (~2015)", "~$1.4–1.6 billion"],
    ["Market Cap (~2015)", "~$6–8 billion"],
    ["Financial Health", "Solidly profitable, strong cash reserves, no debt concerns"],
    ["Pipeline", "Xenotransplantation (genetically modified pig organs)"],
    ["Current Market Cap (2025)", "~$14+ billion (driven by Tyvaso + xenotransplantation)"],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), data,
          col_widths=[Inches(3.0), Inches(9.3)], header_bg=UT_TABLE_HDR, alt_bg=UT_TABLE_ALT)

add_bullets(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(1.5), [
    "Unituxin was a small, complementary diversification into oncology — NOT a transformational asset",
    "Commercial rights obtained from NCI under a CRADA (Cooperative Research & Development Agreement)",
    "Government-to-private technology transfer — not a traditional M&A deal",
    "NCI license terms: upfront fees, milestone payments, royalties on net sales back to NIH/NCI",
], font_size=13, color=MED_GRAY)


# --- SLIDE 12: Unituxin — Revenue & Commercial ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_title_bar(slide, "UNITUXIN  |  Revenue & Commercial Performance",
                  "Small orphan product in a niche market", UT_PRIMARY, UT_ACCENT, UT_SUBTITLE)

data = [
    ["Year", "Estimated Revenue", "Notes"],
    ["2015", "~$12–20M", "Partial year (approved March)"],
    ["2016", "~$30–40M", "Early commercial ramp"],
    ["2017", "~$50–60M", "+$13.5M YoY growth; approaching peak range"],
    ["2018", "~$50–60M", "Stable; peak annual level"],
    ["2024", "Still marketed", "Revenue growth from price increases"],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), data,
          col_widths=[Inches(1.8), Inches(2.5), Inches(8.0)], header_bg=UT_TABLE_HDR, alt_bg=UT_TABLE_ALT)

add_bullets(slide, Inches(0.5), Inches(4.5), Inches(12), Inches(2.5), [
    "Launch price: ~$175,000 per course of treatment",
    "Addressable population: ~700–800 new high-risk neuroblastoma patients/year in U.S., only a subset eligible",
    "Peak annual U.S. sales: ~$50–60M — small by pharma standards but significant for neuroblastoma community",
    "Competition from Danyelza (approved 2020) created additional commercial pressure",
    "PRV voucher (historically valued at $67–350M) potentially had more standalone financial value than Unituxin sales",
], font_size=14, color=DARK_GRAY)


# --- SLIDE 13: Unituxin — Valuation Impact ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_title_bar(slide, "UNITUXIN  |  Valuation Impact on United Therapeutics",
                  "Modest impact on an established, profitable company", UT_PRIMARY, UT_ACCENT, UT_SUBTITLE)

add_bullets(slide, Inches(0.5), Inches(1.6), Inches(12), Inches(5.5), [
    "Stock traded ~$160–180/share around FDA approval (March 2015) — no dramatic movement attributable to Unituxin",
    "Market had largely priced in the approval given strong Phase III clinical data and orphan drug status",
    "PAH franchise (Remodulin, Tyvaso, Orenitram) was always the dominant value driver",
    "PRV voucher potentially had more standalone financial value than Unituxin's sales projections",
    "Unituxin never became a material valuation driver for United Therapeutics",
    "Current state (2025): UTHR market cap ~$14+ billion — driven by Tyvaso growth and xenotransplantation pipeline",
    "Unituxin remains a minor, steady contributor to overall revenue",
], font_size=16, color=DARK_GRAY, spacing=14)

# EU sub-section
add_text_box(slide, Inches(0.5), Inches(5.0), Inches(12), Inches(0.4),
             "EU Market", font_size=18, color=UT_PRIMARY, bold=True)
add_bullets(slide, Inches(0.5), Inches(5.4), Inches(12), Inches(1.8), [
    "Unituxin withdrawn from EU in 2017 due to manufacturing difficulties",
    "Gap filled by Qarziba (dinutuximab beta) — Recordati via EUSA Pharma acquisition ($845M, 2021)",
    "Recordati targets peak Qarziba sales at €300–350M",
], font_size=14, color=MED_GRAY)


# --- SLIDE 14: Key Scientists ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_title_bar(slide, "Key Scientists & Institutions",
                  "The researchers behind anti-GD2 immunotherapy in neuroblastoma",
                  RGBColor(0x2A, 0x2A, 0x3A), RGBColor(0x88, 0x88, 0xAA), RGBColor(0xAA, 0xAA, 0xCC))

# Danyelza scientists
add_text_box(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(0.4),
             "DANYELZA", font_size=18, color=DZ_PRIMARY, bold=True)
data_dz = [
    ["Scientist", "Institution", "Role"],
    ["Dr. Nai-Kong Cheung", "Memorial Sloan Kettering", "Developed 3F8 & hu3F8 (naxitamab);\n30+ years of anti-GD2 research"],
]
add_table(slide, Inches(0.5), Inches(2.0), Inches(12.3), data_dz,
          col_widths=[Inches(2.8), Inches(3.5), Inches(6.0)], header_bg=DZ_TABLE_HDR, alt_bg=DZ_TABLE_ALT)

# Unituxin scientists
add_text_box(slide, Inches(0.5), Inches(3.3), Inches(6), Inches(0.4),
             "UNITUXIN", font_size=18, color=UT_PRIMARY, bold=True)
data_ut = [
    ["Scientist", "Institution", "Role"],
    ["Dr. Ralph Reisfeld", "Scripps Research Institute", "Developed original murine 14.18 antibody"],
    ["Dr. Alice Yu", "UCSD / Rady Children's", "Central clinical investigator;\nled pivotal ANBL0032 Phase III trial (NEJM 2010)"],
    ["Dr. Katherine Matthay", "UCSF Benioff Children's", "Senior COG neuroblastoma investigator"],
    ["Dr. John Maris", "CHOP / UPenn", "Leading neuroblastoma biologist"],
    ["Dr. Stephen Gillies", "—", "Antibody-cytokine fusion proteins"],
    ["NCI CTEP Officers", "NCI, Bethesda, MD", "Managed IND and clinical development"],
]
add_table(slide, Inches(0.5), Inches(3.8), Inches(12.3), data_ut,
          col_widths=[Inches(2.8), Inches(3.5), Inches(6.0)], header_bg=UT_TABLE_HDR, alt_bg=UT_TABLE_ALT)


# ================================================================
# SLIDE 15: CLOSING
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, RGBColor(0x0F, 0x1F, 0x33))
rect(slide, Inches(2), Inches(2.8), Inches(9.333), Inches(0.05), DZ_ACCENT)
add_text_box(slide, Inches(1), Inches(3.1), Inches(11.333), Inches(1.0),
             "Thank You",
             font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.3), Inches(11.333), Inches(0.8),
             "Anti-GD2 Monoclonal Antibodies in Neuroblastoma\nDanyelza (Y-mAbs / SERB)  &  Unituxin (United Therapeutics)",
             font_size=18, color=RGBColor(0xAA, 0xCC, 0xEE), alignment=PP_ALIGN.CENTER)
rect(slide, Inches(2), Inches(5.4), Inches(9.333), Inches(0.05), DZ_ACCENT)


# === SAVE ===
output_path = "/Users/dvn/Desktop/WWWWW/mbr/comp analis/Danyelza_vs_Unituxin_Analysis.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
